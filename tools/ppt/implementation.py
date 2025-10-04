import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from tools.utils import SLIDES_NAME, upload_ppt_to_gcs

CURRENT_DIR = os.getcwd()
IMAGE_PATH = f"{CURRENT_DIR}/assets/logo.jpeg"


class PPTTool:
    def __init__(self):
        """
        Initialize the PPTTool class
        """
        self.prs = None

    def initialize_powerpoint_instance(self):
        """
        Initialize a new PowerPoint presentation instance.

        This sets up a fresh Presentation object using python-pptx, which will hold all the slides that are added later.
        """
        self.prs = Presentation()

    def create_first_slide(self, title: str):
        """
        Create the first slide containing a centered image and a title below it.

        Args:
            title (str): The title text that will be displayed directly below the image.

        Returns:
            None
        """
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)

        image_size_in = Inches(5.2)
        left = (self.prs.slide_width - image_size_in) / 2
        top = (self.prs.slide_height - image_size_in) / 2
        picture = slide.shapes.add_picture(
            IMAGE_PATH, left, top, width=image_size_in, height=image_size_in
        )

        subtitle_top = top + image_size_in
        textbox = slide.shapes.add_textbox(
            left, subtitle_top, image_size_in, Inches(0.6)
        )
        text_frame = textbox.text_frame
        text_frame.text = title

        p = text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.runs[0].font.size = Pt(24)

    def add_text_slide(self, title: str, content: list):
        """
        Create a standard text slide with a title, and content.

        Args:
            title (str): The title string that appears at the top of the slide.
            content (list): A list of content to add.

        Returns:
            None
        """
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[1])
        slide.shapes.title.text = title
        body = slide.shapes.placeholders[1].text_frame

        if isinstance(content, str):
            body.text = content
        elif isinstance(content, list):
            body.text = content[0]
            for item in content[1:]:
                p = body.add_paragraph()
                p.text = item
                p.level = 0
        else:
            raise ValueError("Content must be a string or list of strings")

    def save_slide(self) -> str:
        """
        Save the PowerPoint presentation to a file and uploads it to google cloud.

        Returns:
            str, the URL to the slides uploaded to google cloud.
        """
        try:
            file_path = f"{CURRENT_DIR}/assets/slides/{SLIDES_NAME}.pptx"
            self.prs.save(file_path)
            if os.path.isfile(file_path):
                result = upload_ppt_to_gcs(file_path)
            return result
        except Exception as ex:
            print("====", ex)
